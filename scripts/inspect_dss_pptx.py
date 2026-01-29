"""
Inspect Digital_Security_Services_DSS_Presentation.pptx structure.
Run from project root: python scripts/inspect_dss_pptx.py
Writes scripts/dss_structure.txt (UTF-8).
"""
import os
import sys
from pptx import Presentation

DSS_PATH = os.path.join(os.path.dirname(__file__), "..", "Digital_Security_Services_DSS_Presentation.pptx")
OUT_PATH = os.path.join(os.path.dirname(__file__), "dss_structure.txt")

def safe_str(s, maxlen=200):
    return (s or "")[:maxlen].encode("ascii", "replace").decode("ascii")

def get_text(shape):
    if hasattr(shape, "text"):
        return (shape.text or "").strip()
    if hasattr(shape, "text_frame"):
        return "\n".join((p.text or "") for p in shape.text_frame.paragraphs).strip()
    return ""

def main():
    if not os.path.exists(DSS_PATH):
        sys.stderr.write("DSS file not found: " + DSS_PATH + "\n")
        return
    prs = Presentation(DSS_PATH)
    lines = []
    lines.append("=== DSS Presentation Structure ===\n")
    lines.append("Slide size: {} x {} emu (10 x 7.5 in)\n".format(prs.slide_width, prs.slide_height))
    lines.append("Slides: {}\n".format(len(prs.slides)))
    lines.append("\n--- Slide-by-slide ---\n")
    for i, slide in enumerate(prs.slides, 1):
        lines.append("--- Slide {} ---\n".format(i))
        for sh in slide.shapes:
            t = get_text(sh)
            if t:
                lines.append("  [{}] {}\n".format(sh.shape_type, t[:200] + ("..." if len(t) > 200 else "")))
        lines.append("\n")
    lines.append("=== Layouts ===\n")
    for i, layout in enumerate(prs.slide_layouts):
        lines.append("  Layout {}: {}\n".format(i, getattr(layout, "name", "?")))
    with open(OUT_PATH, "w", encoding="utf-8") as f:
        f.writelines(lines)
    print("Written:", OUT_PATH)
    print("Slides:", len(prs.slides), "| Size: 10 x 7.5 in")

if __name__ == "__main__":
    main()
