"""
Generate SKIFIN presentation from website content.
Structure mirrors Digital_Security_Services_DSS_Presentation.pptx (16 slides):
  Title -> Overview (2 pillars + mission) -> Practice slide -> Section overview -> 10 offering slides -> Summary -> Thank You.
Run: pip install -r requirements-pptx.txt && python generate_skifin_presentation.py
Output: ../SKIFIN_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Output path (project root)
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "..", "SKIFIN_Presentation.pptx")

# Brand colors (dark theme consistent with site)
DARK_BG = RGBColor(0x0D, 0x0D, 0x12)
ACCENT = RGBColor(0x63, 0x66, 0xF1)   # indigo
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA1, 0xA1, 0xAA)


def set_slide_bg(slide, color):
    """Set slide background to solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle=""):
    """DSS-style: Title + subtitle (Slide 1)."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, DARK_BG)
    left, top, width, height = Inches(0.5), Inches(2), Inches(9), Inches(1.2)
    tf = slide.shapes.add_textbox(left, top, width, height)
    p = tf.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        left, top, width, height = Inches(0.5), Inches(3.2), Inches(9), Inches(1)
        tf2 = slide.shapes.add_textbox(left, top, width, height)
        p2 = tf2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = GRAY
    return slide


def add_overview_slide(prs, title, subtitle, pillar1, pillar2, mission):
    """DSS-style Slide 2: Title, subtitle, two pillars, mission statement."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, DARK_BG)
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    p = tf.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    tf2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.5))
    tf2.text_frame.paragraphs[0].text = subtitle
    tf2.text_frame.paragraphs[0].font.size = Pt(16)
    tf2.text_frame.paragraphs[0].font.color.rgb = GRAY
    # Two pillars (left/right)
    for i, (label, left_in) in enumerate([(pillar1, 1.0), (pillar2, 5.0)]):
        box = slide.shapes.add_textbox(Inches(left_in), Inches(2.2), Inches(3.5), Inches(0.8))
        box.text_frame.paragraphs[0].text = label
        box.text_frame.paragraphs[0].font.size = Pt(18)
        box.text_frame.paragraphs[0].font.bold = True
        box.text_frame.paragraphs[0].font.color.rgb = WHITE
    # Mission
    mbox = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(2))
    mbox.text_frame.word_wrap = True
    mbox.text_frame.paragraphs[0].text = "Mission: " + mission
    mbox.text_frame.paragraphs[0].font.size = Pt(14)
    mbox.text_frame.paragraphs[0].font.color.rgb = WHITE
    mbox.text_frame.paragraphs[0].font.italic = True
    return slide


def add_section_title_only(prs, title):
    """DSS-style section slide: single big title, no body (Slides 5–14)."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, DARK_BG)
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(9), Inches(1.4))
    p = tf.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    return slide


def add_practice_slide(prs, title, subtitle, bullets, footer=None):
    """DSS-style Slide 3: Title, subtitle, bullet list, optional footer."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, DARK_BG)
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tf.text_frame.paragraphs[0].text = title
    tf.text_frame.paragraphs[0].font.size = Pt(26)
    tf.text_frame.paragraphs[0].font.bold = True
    tf.text_frame.paragraphs[0].font.color.rgb = ACCENT
    tf2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(9), Inches(0.4))
    tf2.text_frame.paragraphs[0].text = subtitle
    tf2.text_frame.paragraphs[0].font.size = Pt(14)
    tf2.text_frame.paragraphs[0].font.color.rgb = GRAY
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    text_frame = box.text_frame
    text_frame.word_wrap = True
    for i, item in enumerate(bullets):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE
        p.space_after = Pt(6)
    if footer:
        pf = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.5))
        pf.text_frame.paragraphs[0].text = footer
        pf.text_frame.paragraphs[0].font.size = Pt(11)
        pf.text_frame.paragraphs[0].font.color.rgb = GRAY
        pf.text_frame.paragraphs[0].font.italic = True
    return slide


def add_content_slide(prs, title, bullets, sub_bullets=None):
    """Add a content slide with title and bullet list."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, DARK_BG)
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    p = tf.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    text_frame = box.text_frame
    text_frame.word_wrap = True
    for i, item in enumerate(bullets):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)
    if sub_bullets:
        for item in sub_bullets:
            p = text_frame.add_paragraph()
            p.text = "  ◦ " + item
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY
            p.space_after = Pt(4)
    return slide


def add_section_slide(prs, section_title):
    """Add a section divider slide."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, ACCENT)
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    p = tf.text_frame.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    return slide


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add slide with two columns of bullets."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, DARK_BG)
    # Title
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf.text_frame.paragraphs[0].text = title
    tf.text_frame.paragraphs[0].font.size = Pt(26)
    tf.text_frame.paragraphs[0].font.bold = True
    tf.text_frame.paragraphs[0].font.color.rgb = ACCENT
    # Left column
    lbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(5.5))
    lf = lbox.text_frame
    p = lf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    for item in left_items:
        p = lf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_after = Pt(4)
    # Right column
    rbox = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.2), Inches(5.5))
    rf = rbox.text_frame
    p = rf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    for item in right_items:
        p = rf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_after = Pt(4)
    return slide


def add_stats_slide(prs, title, stats):
    """Add slide with stat cards (value + label pairs)."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, DARK_BG)
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf.text_frame.paragraphs[0].text = title
    tf.text_frame.paragraphs[0].font.size = Pt(26)
    tf.text_frame.paragraphs[0].font.bold = True
    tf.text_frame.paragraphs[0].font.color.rgb = ACCENT
    # 2x2 or 1x4 grid of stats
    for i, (value, label) in enumerate(stats):
        col = i % 4
        row = i // 4
        left = Inches(0.6 + col * 2.3)
        top = Inches(1.4 + row * 1.8)
        box = slide.shapes.add_textbox(left, top, Inches(2.1), Inches(1.4))
        tf = box.text_frame
        tf.paragraphs[0].text = value
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(12)
        p2.font.color.rgb = GRAY
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # DSS-style flow (16 slides)

    # 1. Title (like DSS Slide 1)
    add_title_slide(
        prs,
        "SKIFIN",
        "AI-Assisted Solutions Through Vibe Coding | Professional AI & IT Services"
    )

    # 2. Overview — two pillars + mission (like DSS Slide 2)
    add_overview_slide(
        prs,
        title="SKIFIN Overview",
        subtitle="Introduction to SKIFIN",
        pillar1="AI Services (AS) — Automation, Security, DevOps, AI/ML, Data, Cloud, Vibe Coding",
        pillar2="Solutions (SO) — Subsman, Intelligent CRM",
        mission="Deliver end-to-end, scalable AI and IT solutions that accelerate digital transformation."
    )

    # 3. Practice slide — AI Services (like DSS Slide 3: Professional Services)
    add_practice_slide(
        prs,
        title="AI Services",
        subtitle="Consulting & Delivery",
        bullets=[
            "Architecture & Solution Design",
            "Product Deployment & Implementation",
            "Proof of Concepts (PoC), Pilots",
            "AI Automation, Cloud Security, DevOps, AI/ML, Data Analytics, Cloud Migration",
            "Vibe Coding — AI-augmented development",
        ],
        footer="Coverage: Across all technologies in the SKIFIN portfolio."
    )

    # 4. Section: Offerings Overview (like DSS Slide 4: Managed Services Overview)
    add_section_title_only(prs, "Offerings Overview")

    # 5–14. One section-title slide per offering (like DSS Slides 5–14; 10 slides)
    offering_titles = [
        "AI Automation",
        "Cloud Security",
        "DevOps & DevSecOps",
        "AI/ML Solutions",
        "Data Analytics",
        "Cloud Migration",
        "Vibe Coding",
        "Implementation & Support",
        "Subsman — Hyper-local Delivery",
        "Intelligent CRM",
    ]
    for title in offering_titles:
        add_section_title_only(prs, title)

    # 15. Summary & Differentiators (like DSS Slide 15)
    add_content_slide(prs, "Summary & Differentiators", [
        "500+ projects delivered | 98% client satisfaction | 50+ enterprise clients | 24/7 support",
        "15+ years experience | Open-source first | Scalable, secure solutions",
        "AI Services: Agentic AI, RPA, Zero Trust, CI/CD, Generative AI, BI, Cloud Migration, Vibe Coding",
        "Solutions: Subsman (milk, water, tiffin, grocery) and Intelligent CRM (pipelines, analytics)",
        "Mission: End-to-end, secure, scalable solutions across the AI and IT lifecycle.",
    ])

    # 16. Thank You / Contact Info (like DSS Slide 16)
    add_title_slide(
        prs,
        "Thank You / Contact",
        "contact@skifin.com  |  +91-8700605399  |  Galaxy Business Park, Sector 62, Noida"
    )

    prs.save(OUTPUT_PATH)
    print("Presentation saved to:", os.path.abspath(OUTPUT_PATH))


if __name__ == "__main__":
    main()
